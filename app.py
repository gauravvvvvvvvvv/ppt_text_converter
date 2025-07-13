import streamlit as st
from pptx import Presentation
from pptx.table import Table
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from io import BytesIO
import os
import zipfile
import shutil
import xml.etree.ElementTree as ET
import tempfile

st.set_page_config(page_title="Balaram to Unicode Converter", page_icon="📘", layout="centered")

def load_css():
    st.markdown("""<style>/* your full CSS here */</style>""", unsafe_allow_html=True)

load_css()

st.markdown("<h1>📘 Balaram to Unicode PPTX Converter</h1>", unsafe_allow_html=True)

uploaded_file = st.file_uploader("📂 Upload your .pptx file", type=["pptx"])
just_unlock = st.checkbox("🔓 Only unlock file (no conversion)", value=False)

balaram_map = {
    'ä': 'ā', 'é': 'ī', 'ü': 'ū', 'å': 'ṛ', 'è': 'ṝ', 'ì': 'ṅ', 'ï': 'ñ', 'ö': 'ṭ',
    'ò': 'ḍ', 'ë': 'ṇ', 'ç': 'ś', 'à': 'ṁ', 'ù': 'ḥ', 'ÿ': 'ḷ', 'û': 'ḹ', 'ý': 'ẏ',
    'Ä': 'Ā', 'É': 'Ī', 'Ü': 'Ū', 'Å': 'Ṛ', 'È': 'Ṝ', 'Ì': 'Ṅ', 'Ï': 'Ñ',
    'Ö': 'Ṭ', 'Ò': 'Ḍ', 'Ë': 'Ṇ', 'Ç': 'Ś', 'À': 'Ṁ', 'Ù': 'Ḥ', 'ß': 'Ḷ',
    'Ý': 'Ẏ', '~': 'ɱ', "'": "'", '…': '…', '’': '’', 'ñ': 'ṣ', 'Ñ': 'Ṣ'
}

def convert_balaram_to_unicode(text: str) -> str:
    return ''.join(balaram_map.get(char, char) for char in text)

def convert_text_frame(tf):
    if tf and tf.text.strip():
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = convert_balaram_to_unicode(run.text)
        return True
    return False

def convert_table(table: Table):
    conversions = 0
    for row in table.rows:
        for cell in row.cells:
            if convert_text_frame(cell.text_frame):
                conversions += 1
    return conversions

def process_shape(shape):
    conversions = 0
    if isinstance(shape, Picture): return 0
    try:
        if shape.has_text_frame:
            if convert_text_frame(shape.text_frame):
                conversions += 1
        elif hasattr(shape, 'shape_type') and shape.shape_type == 19:
            conversions += convert_table(shape.table)
        elif isinstance(shape, GroupShape):
            for subshape in shape.shapes:
                conversions += process_shape(subshape)
    except:
        pass
    return conversions

def unlock_pptx_file(pptx_bytes, filename):
    with tempfile.TemporaryDirectory() as temp_dir:
        zip_temp = os.path.join(temp_dir, "temp.zip")
        extract_path = os.path.join(temp_dir, "extract")
        with open(zip_temp, 'wb') as f:
            f.write(pptx_bytes)
        try:
            with zipfile.ZipFile(zip_temp, 'r') as zip_ref:
                zip_ref.extractall(extract_path)
        except:
            return pptx_bytes
        pres_xml = os.path.join(extract_path, 'ppt', 'presentation.xml')
        if os.path.exists(pres_xml):
            try:
                ET.register_namespace('p', "http://schemas.openxmlformats.org/presentationml/2006/main")
                tree = ET.parse(pres_xml)
                root = tree.getroot()
                for elem in root.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}modifyVerifier'):
                    root.remove(elem)
                for elem in root.findall('.//modifyVerifier'):
                    try:
                        root.remove(elem)
                    except:
                        pass
                tree.write(pres_xml, encoding='utf-8', xml_declaration=True)
            except:
                pass
        try:
            output_zip = os.path.join(temp_dir, "unlocked.zip")
            shutil.make_archive(output_zip.replace('.zip', ''), 'zip', extract_path)
            with open(output_zip, 'rb') as f:
                return f.read()
        except:
            return pptx_bytes

def convert_pptx(pptx_bytes):
    try:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_input:
            tmp_input.write(pptx_bytes)
            tmp_input.flush()
            tmp_path = tmp_input.name

        prs = Presentation(tmp_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                process_shape(shape)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_output:
            prs.save(tmp_output.name)
            tmp_output.seek(0)
            with open(tmp_output.name, 'rb') as f:
                return BytesIO(f.read())

    except Exception as e:
        print(f"❌ Error during conversion: {e}")
        return None

if uploaded_file is not None:
    try:
        file_bytes = uploaded_file.getvalue()
        st.write(f"📄 File size: `{len(file_bytes) / 1024**2:.2f} MB`")
        unlocked_bytes = unlock_pptx_file(file_bytes, uploaded_file.name)
        st.write(f"🔓 Unlocked size: `{len(unlocked_bytes) / 1024**2:.2f} MB`")

        if just_unlock:
            st.download_button(
                label="📥 Download Unlocked PPTX",
                data=unlocked_bytes,
                file_name=f"{os.path.splitext(uploaded_file.name)[0]}_unlocked.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
        else:
            converted_stream = convert_pptx(unlocked_bytes)
            if converted_stream:
                st.download_button(
                    label="📥 Download Converted PPTX",
                    data=converted_stream,
                    file_name=f"{os.path.splitext(uploaded_file.name)[0]}_unicode.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
            else:
                st.error("❌ Conversion failed.")
    except Exception as e:
        st.error(f"❌ Error: {e}")

with st.expander("ℹ️ How to use this converter"):
    st.markdown("""
    1. Upload a `.pptx` file  
    2. Choose whether to unlock only or convert  
    3. Download the result  
    """)

st.markdown(
    "<div style='text-align:center; font-size:16px; margin-top: 20px; color: #6d3600;'>"
    "🌸 Hare Kṛṣṇa! All glories to Śrīla Prabhupāda. 🌸</div>",
    unsafe_allow_html=True
)
